"""PPTX extractor — one segment per slide, via python-pptx."""
from __future__ import annotations
from pathlib import Path
from pptx import Presentation
from semanticsd.extractors.base import Extractor, ExtractedDoc, ExtractedSegment
from semanticsd.extractors.registry import register


@register
class PptxExtractor(Extractor):
    file_type = "pptx"
    extensions = (".pptx",)

    def extract(self, path: Path) -> ExtractedDoc:
        prs = Presentation(str(path))
        segments: list[ExtractedSegment] = []
        cursor = 0
        for i, slide in enumerate(prs.slides, start=1):
            parts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    parts.append(shape.text.strip())
            text = "\n".join(p for p in parts if p)
            if not text:
                continue
            seg_bytes = len(text.encode("utf-8"))
            segments.append(ExtractedSegment(
                text=text,
                byte_start=cursor,
                byte_end=cursor + seg_bytes,
                metadata={"slide": i},
            ))
            cursor += seg_bytes + 1
        return ExtractedDoc(
            path=str(path),
            file_type=self.file_type,
            segments=segments,
        )
