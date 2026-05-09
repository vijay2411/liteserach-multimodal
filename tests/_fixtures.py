"""Programmatic fixture builders for extractor tests.

Each helper returns a Path to a tmp file containing a minimal valid file of
that type. The caller (a test using the `tmp_path` pytest fixture) supplies
the parent dir, so files are auto-cleaned.
"""
from __future__ import annotations
from pathlib import Path


def make_text(tmp_path: Path, name: str = "doc.txt", body: str = "Hello world.\nLine two.") -> Path:
    p = tmp_path / name
    p.write_text(body)
    return p


def make_markdown(tmp_path: Path, name: str = "doc.md") -> Path:
    return make_text(tmp_path, name=name, body="# Title\n\nA paragraph.\n\n- bullet\n")


def make_code(tmp_path: Path, name: str = "hello.py") -> Path:
    return make_text(tmp_path, name=name, body='def greet():\n    return "hi"\n')


def make_json(tmp_path: Path, name: str = "data.json") -> Path:
    return make_text(tmp_path, name=name, body='{"name": "alice", "tags": ["a", "b"]}')


def make_yaml(tmp_path: Path, name: str = "data.yaml") -> Path:
    return make_text(tmp_path, name=name, body="name: alice\ntags:\n  - a\n  - b\n")


def make_csv(tmp_path: Path, name: str = "data.csv") -> Path:
    return make_text(tmp_path, name=name, body="name,tag\nalice,a\nbob,b\n")


def make_html(tmp_path: Path, name: str = "page.html") -> Path:
    body = (
        "<!DOCTYPE html><html><head><title>Page</title></head>"
        "<body><h1>Heading</h1><p>Paragraph one.</p>"
        "<script>console.log('skip')</script></body></html>"
    )
    return make_text(tmp_path, name=name, body=body)


def make_pdf(tmp_path: Path, name: str = "doc.pdf") -> Path:
    """Two-page PDF with simple text, generated via reportlab."""
    from reportlab.pdfgen import canvas
    p = tmp_path / name
    c = canvas.Canvas(str(p))
    c.drawString(100, 700, "Hello from page one.")
    c.showPage()
    c.drawString(100, 700, "Page two has different text.")
    c.showPage()
    c.save()
    return p


def make_docx(tmp_path: Path, name: str = "doc.docx") -> Path:
    from docx import Document
    p = tmp_path / name
    doc = Document()
    doc.add_heading("Title", level=1)
    doc.add_paragraph("First paragraph.")
    doc.add_paragraph("Second paragraph with different content.")
    doc.save(str(p))
    return p


def make_xlsx(tmp_path: Path, name: str = "book.xlsx") -> Path:
    from openpyxl import Workbook
    p = tmp_path / name
    wb = Workbook()
    ws = wb.active
    ws.title = "data"
    ws.append(["name", "score"])
    ws.append(["alice", 90])
    ws.append(["bob", 85])
    wb.save(str(p))
    return p


def make_pptx(tmp_path: Path, name: str = "deck.pptx") -> Path:
    from pptx import Presentation
    p = tmp_path / name
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]  # Title only
    s = prs.slides.add_slide(slide_layout)
    s.shapes.title.text = "Slide One"
    s2 = prs.slides.add_slide(slide_layout)
    s2.shapes.title.text = "Slide Two has different content"
    prs.save(str(p))
    return p


def make_epub(tmp_path: Path, name: str = "book.epub") -> Path:
    from ebooklib import epub
    p = tmp_path / name
    book = epub.EpubBook()
    book.set_identifier("id-1")
    book.set_title("Tiny Book")
    book.set_language("en")
    c1 = epub.EpubHtml(title="Chapter 1", file_name="c1.xhtml", lang="en")
    c1.content = "<h1>Chapter 1</h1><p>Lorem ipsum content.</p>"
    book.add_item(c1)
    book.toc = (epub.Link("c1.xhtml", "Chapter 1", "c1"),)
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())
    book.spine = ["nav", c1]
    epub.write_epub(str(p), book)
    return p


def make_rtf(tmp_path: Path, name: str = "doc.rtf") -> Path:
    body = r"{\rtf1\ansi\deff0 {\fonttbl{\f0 Helvetica;}}\f0\fs24 Hello from RTF.\par Second line.\par}"
    return make_text(tmp_path, name=name, body=body)


def make_eml(tmp_path: Path, name: str = "msg.eml") -> Path:
    body = (
        "From: alice@example.com\r\n"
        "To: bob@example.com\r\n"
        "Subject: A test message\r\n"
        "Content-Type: text/plain; charset=utf-8\r\n"
        "\r\n"
        "Hello Bob, this is the body of the message.\r\n"
    )
    return make_text(tmp_path, name=name, body=body)


def make_ipynb(tmp_path: Path, name: str = "notebook.ipynb") -> Path:
    import json
    nb = {
        "nbformat": 4,
        "nbformat_minor": 5,
        "metadata": {},
        "cells": [
            {"cell_type": "markdown", "metadata": {}, "source": ["# Notebook Title\n", "Some markdown text."]},
            {"cell_type": "code", "metadata": {}, "execution_count": None, "outputs": [],
             "source": ["x = 42\n", "print(x)"]},
        ],
    }
    p = tmp_path / name
    p.write_text(json.dumps(nb))
    return p


def make_image_with_text(tmp_path: Path, name: str = "img.png", text: str = "hello world") -> Path:
    """A PNG with rendered text — used for OCR tests.

    Uses a TrueType font when available (better Tesseract accuracy) and falls
    back to Pillow's default bitmap font otherwise. The image is intentionally
    large to give Tesseract enough resolution.
    """
    from PIL import Image, ImageDraw, ImageFont
    p = tmp_path / name
    img = Image.new("RGB", (600, 200), color="white")
    draw = ImageDraw.Draw(img)
    font = None
    for candidate in (
        "/System/Library/Fonts/Helvetica.ttc",
        "/System/Library/Fonts/Supplemental/Arial.ttf",
    ):
        try:
            font = ImageFont.truetype(candidate, 60)
            break
        except (OSError, IOError):
            continue
    if font is None:
        font = ImageFont.load_default()
    draw.text((20, 60), text, fill="black", font=font)
    img.save(str(p))
    return p


def make_wav_silence(tmp_path: Path, name: str = "audio.wav", duration_s: float = 0.5) -> Path:
    """A tiny silent WAV. faster-whisper transcribes it to '' or near-empty."""
    import wave
    p = tmp_path / name
    with wave.open(str(p), "wb") as w:
        w.setnchannels(1)
        w.setsampwidth(2)
        w.setframerate(16000)
        n_frames = int(16000 * duration_s)
        w.writeframes(b"\x00\x00" * n_frames)
    return p
